def scanner_node(state: AgentState):
    prs = Presentation(state["template_path"])
    guide_lines = ["=== [통합] 템플릿 선택 가이드 ==="]
    
    # 레이아웃 정보 저장용 (Planner가 쓸 인덱스 매핑)
    layout_map = {} 

    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        layout_map[layout_name] = i
        
        # [A] Dynamic 레이아웃일 때 (이름이 Dynamic_으로 시작)
        if layout_name.startswith("Dynamic_"):
            info = f"\n[Layout Index: {i}] 타입: 🔧Dynamic (차트/표/자유배치용) | 이름: {layout_name}"
            info += "\n   👉 사용 가능한 가이드(Anchor):"
            
            # 가이드 도형 찾기 (Guide_로 시작하는 도형)
            anchors = [s.name for s in layout.shapes if s.name.startswith("Guide_")]
            if anchors:
                info += f" {', '.join(anchors)}"
            else:
                info += " (가이드 도형 없음)"
            guide_lines.append(info)

        # [B] Static 레이아웃일 때 (기존 방식)
        else:
            info = f"\n[Layout Index: {i}] 타입: 📄Static (정형 텍스트/이미지용) | 이름: {layout_name}"
            info += "\n   👉 채워야 할 칸(Placeholder):"
            
            placeholders = [s.name for s in layout.placeholders]
            info += f" {', '.join(placeholders)}"
            guide_lines.append(info)

    return {"template_guide": "\n".join(guide_lines)}



def planner_node(state: AgentState):
    guide = state["template_guide"]
    
    system_prompt = """
    당신은 PPT 스토리보드 작가입니다. 사용자 요청을 분석하여 **논리적인 흐름을 갖춘 여러 장의 슬라이드**를 기획하세요.
    
    [작성 전략]
    1. **표지/목차/간지** 등 정형화된 페이지는 -> **'static'** 타입 사용.
    2. **데이터 시각화(차트, 복잡한 표)**가 필요한 페이지는 -> **'dynamic'** 타입 사용.
    
    [응답 형식: JSON List]
    [
        {
            "type": "static",
            "layout_index": 0,
            "content_mapping": { "Title": "전기차 시장 분석", "Subtitle": "2024 Report" }
        },
        {
            "type": "dynamic",
            "layout_index": 5,
            "title": "시장 점유율 현황",
            "components": [
                { "type": "chart", "position": "Guide_Left", "data": {...} },
                { "type": "text", "position": "Guide_Right", "content": "..." }
            ]
        }
    ]

    [템플릿 가이드]
    {guide}
    """
    
    # ... (LLM 호출 및 JSON 파싱 로직은 이전과 동일) ...
    # 결과로 List[dict] 형태의 slide_data를 반환합니다.



def renderer_node(state: AgentState):
    print("--- [Node 3] 통합 렌더링 시작 ---")
    slides_data = state["slide_data"] # 리스트
    prs = Presentation(state["template_path"])
    
    for plan in slides_data:
        layout_idx = plan["layout_index"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # [모드 1] Static (기존 채우기 방식)
        if plan["type"] == "static":
            print(f"📄 Static 슬라이드 생성: Layout {layout_idx}")
            mapping = plan["content_mapping"]
            
            for shape in slide.placeholders:
                if shape.name in mapping:
                    content = mapping[shape.name]
                    # (기존의 텍스트/이미지 삽입 함수 호출)
                    # insert_text(shape, content) or insert_image(...)

        # [모드 2] Dynamic (앵커 기반 그리기 방식)
        elif plan["type"] == "dynamic":
            print(f"🔧 Dynamic 슬라이드 생성: Layout {layout_idx}")
            
            # 1. 제목 설정 (제목 Placeholder는 보통 공통적으로 존재하므로 처리)
            if slide.shapes.title:
                slide.shapes.title.text = plan.get("title", "")
            
            # 2. 앵커(Guide) 도형 위치 파악
            anchors = {}
            for shape in slide.shapes:
                if shape.name.startswith("Guide_"):
                    anchors[shape.name] = (shape.left, shape.top, shape.width, shape.height)
                    # (선택) 가이드 도형 숨기기: shape.visible = False
            
            # 3. 컴포넌트 그리기
            for comp in plan["components"]:
                pos_name = comp["position"]
                if pos_name in anchors:
                    x, y, w, h = anchors[pos_name]
                    
                    if comp["type"] == "chart":
                        draw_chart(slide, x, y, w, h, comp["data"])
                    elif comp["type"] == "table":
                        draw_table(slide, x, y, w, h, comp["data"])
                    elif comp["type"] == "text":
                        draw_text(slide, x, y, w, h, comp["content"])
                else:
                    print(f"⚠️ 앵커 '{pos_name}'를 찾을 수 없음")

    prs.save(state["output_path"])
    return {"final_message": "완료"}


# template_config.py

TEMPLATE_REGISTRY = {
    # === [Static] 정형화된 템플릿 ===
    "Title_Slide": {
        "type": "static",
        "desc": "프레젠테이션의 표지입니다. 제목과 부제목만 들어갑니다.",
        "rules": {"Title": "20자 이내", "Subtitle": "날짜/발표자 포함"}
    },
    "Agenda_Slide": {
        "type": "static",
        "desc": "목차를 나열할 때 사용합니다.",
        "rules": {"Content": "개조식으로 작성"}
    },

    # === [Dynamic] 자유 배치 템플릿 ===
    "Dynamic_Split": {
        "type": "dynamic",
        "desc": "두 가지 항목(예: 매출 비교, 경쟁사 분석)을 좌우로 비교할 때 씁니다.",
        "anchors": ["Guide_Left", "Guide_Right"] # (참고용: 실제 파일과 일치해야 함)
    },
    "Dynamic_Full": {
        "type": "dynamic",
        "desc": "복잡한 대형 표나 차트 하나를 크게 보여줄 때 씁니다.",
        "anchors": ["Guide_Main"]
    }
}

def generate_template_guide(pptx_path):
    prs = Presentation(pptx_path)
    guide_lines = []
    
    # 1. Static과 Dynamic을 구분해서 보여주기 위해 리스트 분리
    static_guides = ["--- [Static Layouts: 채우기 모드] ---"]
    dynamic_guides = ["--- [Dynamic Layouts: 그리기 모드] ---"]

    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name
        
        # [필터링] 레지스트리에 없는 레이아웃은 LLM에게 안 보여줌 (토큰 절약)
        if name not in TEMPLATE_REGISTRY:
            continue
            
        config = TEMPLATE_REGISTRY[name]
        desc = config.get("desc", "")
        
        # === A. Static 처리 ===
        if config["type"] == "static":
            placeholders = [p.name for p in layout.placeholders]
            info = f"\nCreate [Index: {i}] Name: '{name}'"
            info += f"\n   - 용도: {desc}"
            info += f"\n   - 입력칸: {', '.join(placeholders)}"
            static_guides.append(info)
            
        # === B. Dynamic 처리 ===
        elif config["type"] == "dynamic":
            # 실제 파일에서 'Guide_'로 시작하는 도형 찾기
            real_anchors = [s.name for s in layout.shapes if s.name.startswith("Guide_")]
            
            info = f"\nCreate [Index: {i}] Name: '{name}'"
            info += f"\n   - 용도: {desc}"
            info += f"\n   - 가이드 영역(Anchors): {', '.join(real_anchors)}"
            dynamic_guides.append(info)

    # 두 그룹을 합쳐서 리턴
    full_text = "\n".join(static_guides + ["\n"] + dynamic_guides)
    return full_text


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# --- [도구 함수들] ---
def draw_chart(slide, x, y, w, h, data):
    # data format: {"labels": ["A", "B"], "values": [10, 20]}
    chart_data = CategoryChartData()
    chart_data.categories = data.get('labels', [])
    chart_data.add_series('Series 1', data.get('values', []))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data)

def draw_table(slide, x, y, w, h, data_rows):
    rows = len(data_rows)
    cols = len(data_rows[0])
    graphic_frame = slide.shapes.add_table(rows, cols, x, y, w, h)
    table = graphic_frame.table
    for r in range(rows):
        for c in range(cols):
            table.cell(r, c).text = str(data_rows[r][c])

def draw_text(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.text = text
    # 스타일링 (Dynamic 모드는 폰트 지정 필요)
    for p in tb.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.name = "AppleSDGothicNeo"

# --- [Renderer 메인 로직] ---
def renderer_node(state: AgentState):
    prs = Presentation(state["template_path"])
    
    # 캔버스 작업 영역 정의 (제목, 로고 제외한 빈 공간)
    canvas_x = Inches(0.5)
    canvas_y = Inches(1.5) # 제목 아래부터 시작
    canvas_w = Inches(9.0)
    canvas_h = Inches(5.0)

    for slide_plan in state["slide_data"]:
        # 1. 템플릿 모드 (기존 로직)
        if slide_plan["type"] == "template":
            # ... (기존 코드) ...
            pass
            
        # 2. 다이내믹 모드 (신규 로직)
        elif slide_plan["type"] == "dynamic":
            slide = prs.slides.add_slide(prs.slide_layouts[slide_plan["layout_index"]])
            
            # 제목은 Placeholder 사용 (일관성 유지)
            if slide.shapes.title:
                slide.shapes.title.text = slide_plan.get("title", "")
            
            # 레이아웃 계산 (간단한 그리드 시스템)
            plan = slide_plan["layout_plan"]
            comps = slide_plan["components"]
            
            # 영역 분할 로직
            regions = []
            if plan == "Split_Left_Right":
                regions = [
                    (canvas_x, canvas_y, canvas_w/2 - Inches(0.2), canvas_h), # 왼쪽
                    (canvas_x + canvas_w/2 + Inches(0.2), canvas_y, canvas_w/2 - Inches(0.2), canvas_h) # 오른쪽
                ]
            else: # Full
                regions = [(canvas_x, canvas_y, canvas_w, canvas_h)]
            
            # 컴포넌트 그리기
            for i, comp in enumerate(comps):
                if i >= len(regions): break # 영역보다 컴포넌트가 많으면 무시
                x, y, w, h = regions[i]
                
                if comp["type"] == "chart":
                    draw_chart(slide, x, y, w, h, comp["data"])
                elif comp["type"] == "table":
                    draw_table(slide, x, y, w, h, comp["data"])
                elif comp["type"] == "text":
                    draw_text(slide, x, y, w, h, comp["content"])
                elif comp["type"] == "image":
                     # 이미지 파일명으로 바이너리 찾아서 삽입 (이전 코드 활용)
                     pass

    prs.save(state["output_path"])


def get_real_placeholder_name(slide_shape):
    """
    슬라이드에 있는 쉐이프(Placeholder)의 진짜 이름을 
    원본 레이아웃(Master)에서 찾아옵니다.
    """
    try:
        # 1. 현재 상자의 고유 ID (idx) 확인
        idx = slide_shape.placeholder_format.idx
        
        # 2. 이 슬라이드를 만든 '부모 레이아웃'을 호출
        layout = slide_shape.part.slide_layout
        
        # 3. 레이아웃에서 똑같은 idx를 가진 상자를 찾음
        original_shape = layout.placeholders[idx]
        
        # 4. 그 상자의 이름(우리가 지어준 이름)을 반환
        return original_shape.name
        
    except Exception as e:
        # 일반 도형이거나 찾을 수 없으면 그냥 현재 이름 반환
        return slide_shape.name


from pptx import Presentation

def inspect_template(pptx_path):
    prs = Presentation(pptx_path)
    print(f"=== 템플릿 분석: {pptx_path} ===\n")

    for i, layout in enumerate(prs.slide_layouts):
        print(f"📄 [Layout {i}] 이름: {layout.name}")
        
        # 레이아웃에 있는 모든 Placeholder 조회
        for ph in layout.placeholders:
            # idx: 주민등록번호 (변하지 않음)
            # name: 우리가 지어준 이름
            # type: 텍스트인지, 이미지인지, 제목인지
            print(f"   - [idx: {ph.placeholder_format.idx}] 이름: '{ph.name}' (Type: {ph.placeholder_format.type})")
        print("-" * 30)

# 실행
inspect_template("my_template.pptx")


def renderer_node(state: AgentState):
    # ... (생략) ...
    
    # 헬퍼 함수: idx로 진짜 이름 찾기 (이전 단계에서 만든 것)
    def get_real_ph_name(slide_ph):
        current_idx = slide_ph.placeholder_format.idx
        for parent_ph in slide_ph.part.slide_layout.placeholders:
            if parent_ph.placeholder_format.idx == current_idx:
                return parent_ph.name
        return slide_ph.name

    for plan in slides_data:
        layout_idx = plan["layout_index"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # ====================================================
        # [PART A] Placeholder 채우기 (Title, Subtitle, Static 내용)
        # ====================================================
        # AI가 준 데이터 중 'content_mapping'과 'common_fields'를 합쳐서 봅니다.
        fill_data = {**plan.get("content_mapping", {}), **plan.get("common_fields", {})}
        
        for shape in slide.placeholders:
            real_name = get_real_ph_name(shape) # idx 기반 역추적
            
            # Title 예외처리
            if shape.placeholder_format.type == 1: 
                real_name = "Title"

            if real_name in fill_data:
                # 텍스트 삽입 로직
                shape.text = fill_data[real_name] 

        # ====================================================
        # [PART B] Anchor 위에 그리기 (Dynamic Components)
        # ====================================================
        if plan.get("type") == "dynamic":
            components = plan.get("components", [])
            
            # 1. 레이아웃 원본에서 앵커 좌표 수집 (슬라이드X -> 레이아웃O)
            layout = prs.slide_layouts[layout_idx]
            anchors = {}
            
            for shape in layout.shapes:
                # Placeholder가 아니고, 이름이 Guide_로 시작하는 것
                if not shape.is_placeholder and shape.name.startswith("Guide_"):
                    anchors[shape.name] = (shape.left, shape.top, shape.width, shape.height)
            
            # 2. 좌표에 맞춰 그리기
            for comp in components:
                target_name = comp["position"]
                
                if target_name in anchors:
                    x, y, w, h = anchors[target_name]
                    
                    # (그리기 함수 호출)
                    if comp["type"] == "chart":
                        draw_chart(slide, x, y, w, h, comp["data"])
                    elif comp["type"] == "table":
                        draw_table(slide, x, y, w, h, comp["data"])
                    elif comp["type"] == "text":
                        # Anchor 위에 텍스트 상자를 새로 그림
                        draw_text_box(slide, x, y, w, h, comp["content"])
                else:
                    print(f"⚠️ 앵커 못 찾음: {target_name}")
                    
    # ...




