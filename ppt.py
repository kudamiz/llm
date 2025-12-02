from pptx import Presentation
import matplotlib.pyplot as plt

# 1. 템플릿 불러오기
prs = Presentation('template.pptx')

# 2. 사용할 레이아웃 선택 (위에서 확인한 Index [1]번 레이아웃 사용)
slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(slide_layout)

# 3. 데이터 준비 (LangGraph에서 받아온 데이터라고 가정)
title_text = "2025 반도체 공정 수율 분석"
body_text = "- 주요 이슈: 식각 공정 온도 변화\n- 조치 사항: 챔버 3호기 파라미터 조정 완료"
# 차트는 이미지로 저장해둡니다
plt.plot([1, 2, 3], [10, 20, 5])
plt.savefig('chart_output.png') 

# 4. 데이터 주입 (Placeholder Index 활용)
# (1) 제목 넣기 (Index 0)
slide.placeholders[0].text = title_text

# (2) 본문 텍스트 넣기 (Index 10)
slide.placeholders[10].text = body_text

# (3) 차트 이미지 넣기 (Index 11)
# 이미지는 placeholder 자리에 쏙 들어갑니다. (자동 리사이징 됨)
slide.placeholders[11].insert_picture('chart_output.png')

# 5. 저장
prs.save('Final_Report.pptx')
print("PPT 생성 완료!")



from pptx import Presentation

def analyze_ppt_layout(ppt_path):
    prs = Presentation(ppt_path)
    
    print(f"파일: {ppt_path} 분석 시작\n")
    
    # 슬라이드 마스터에 있는 레이아웃들을 순회
    for i, layout in enumerate(prs.slide_layouts):
        print(f"--- Layout Index [{i}]: {layout.name} ---")
        
        # 해당 레이아웃 안에 있는 Placeholder(틀)들을 순회
        for shape in layout.placeholders:
            print(f"   Placeholder Index [{shape.placeholder_format.idx}] "
                  f"- 이름: {shape.name} (타입: {shape.placeholder_format.type})")

# 사용 예시
analyze_ppt_layout("template.pptx")
