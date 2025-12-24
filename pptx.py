7from pptx import Presentation

def get_template_guide(pptx_path):
    prs = Presentation(pptx_path)
    guide_text = "현재 사용 가능한 PPT 레이아웃 목록입니다:\n"
    
    # 모든 마스터 레이아웃을 순회
    for i, layout in enumerate(prs.slide_layouts):
        # 레이아웃 이름 (예: Comparison)
        layout_info = {
            "layout_index": i,
            "layout_name": layout.name,
            "placeholders": []
        }
        
        # 레이아웃 안의 구멍(Placeholder)들 이름 수집
        for shape in layout.placeholders:
            # PPT '선택 창'에서 지정한 이름을 그대로 가져옴
            p_info = f"{shape.name} (ID: {shape.placeholder_format.idx})"
            layout_info["placeholders"].append(p_info)
            
        guide_text += str(layout_info) + "\n"
        
    return guide_text

# 실행 결과 예시 (이 텍스트가 자동으로 생성됨)
# "{'layout_index': 1, 'layout_name': '2단비교', 'placeholders': ['Title (ID:0)', 'Body_Left (ID:1)', 'Body_Right (ID:2)']}"


# 시스템 프롬프트 템플릿
system_prompt = """
당신은 PPT 생성 전문가입니다. 
아래 제공된 [템플릿 가이드]를 보고, 사용자 입력에 가장 적합한 layout_index를 선택하고,
각 placeholder 이름에 맞는 내용을 JSON으로 생성하세요.

[템플릿 가이드]
{template_guide}  <-- 여기에 파이썬이 읽은 정보가 자동으로 들어감
"""

# 실행 시점
current_guide = get_template_guide("company_template_v2.pptx") # 파일만 바꾸면 됨
formatted_prompt = system_prompt.format(template_guide=current_guide)

from typing import List, Dict
from pydantic import BaseModel, Field
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate

# 1. LLM이 뱉어내야 할 최종 데이터 구조 정의 (Schema)
class SlideOutput(BaseModel):
    layout_index: int = Field(..., description="선택한 슬라이드 레이아웃의 인덱스 번호")
    # key: placeholder 이름, value: 들어갈 내용
    content_mapping: Dict[str, str] = Field(..., description="Placeholder 이름을 키(Key)로, 채울 내용을 값(Value)으로 하는 딕셔너리")
    reason: str = Field(..., description="이 레이아웃을 선택한 이유")

# 2. 에이전트 함수 정의
def generate_slide_json(user_input: str, template_guide: str):
    # 모델 설정 (JSON 모드 지원하는 모델 권장)
    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    
    # 구조화된 출력을 하도록 설정
    structured_llm = llm.with_structured_output(SlideOutput)

    # 프롬프트 구성 (동적 템플릿 가이드 주입)
    system_prompt = """
    당신은 PPT 생성 전문가입니다. 
    사용자의 입력을 분석하고, 아래 [템플릿 가이드]를 참고하여 가장 적절한 레이아웃을 선택하세요.
    그리고 각 Placeholder의 'Name'에 맞춰 내용을 요약/배치하여 JSON으로 반환하세요.
    
    [템플릿 가이드]
    {guide}
    """
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", system_prompt),
        ("human", "{input}")
    ])

    # 실행 체인
    chain = prompt | structured_llm
    
    # 결과 반환 (Pydantic 객체)
    return chain.invoke({"guide": template_guide, "input": user_input})
 
from pptx import Presentation

# --- [설정] 파일 경로 ---
TEMPLATE_PATH = "my_template.pptx"   # 준비된 템플릿 파일
OUTPUT_PATH = "step3_test_result.pptx" # 결과 파일

# ====================================================
# [가정] Step 2에서 AI가 만들어줬다고 칠 '가짜 데이터'
# ====================================================
# 주의: 아래 딕셔너리의 Key값("Title", "Body_Left" 등)은 
# 반드시 PPT '선택 창'에서 지정한 이름과 똑같아야 합니다!
mock_slide_data = {
    "layout_index": 1,  # 템플릿에서 테스트하고 싶은 레이아웃 번호 (예: 1번)
    "content_mapping": {
        "Title": "Step 3 단독 테스트 성공!", 
        "Body_Left": "여기는 왼쪽 본문 영역입니다.\n데이터가 잘 들어갔나요?",
        "Body_Right": "여기는 오른쪽 본문 영역입니다.\nPython-pptx로 생성되었습니다."
    }
}

# ====================================================
# 3. Renderer 함수 (파일 생성 로직)
# ====================================================
def create_ppt_file_test(data, template_path, output_path):
    print(f"📂 템플릿 여는 중: {template_path}")
    try:
        prs = Presentation(template_path)
    except FileNotFoundError:
        print("❌ 오류: 템플릿 파일을 찾을 수 없습니다.")
        return

    # 1. 레이아웃 선택
    target_index = data["layout_index"]
    try:
        selected_layout = prs.slide_layouts[target_index]
        print(f"🎨 선택된 레이아웃: '{selected_layout.name}' (Index: {target_index})")
    except IndexError:
        print(f"❌ 오류: 레이아웃 번호 {target_index}번은 존재하지 않습니다.")
        return

    # 2. 슬라이드 추가
    slide = prs.slides.add_slide(selected_layout)

    # 3. 데이터 매핑 (핵심!)
    mapping = data["content_mapping"]
    
    matched_count = 0
    print("\n--- [데이터 매핑 시작] ---")
    
    for shape in slide.placeholders:
        shape_name = shape.name # PPT에 설정된 이름
        
        # 1) 이름이 매칭되는지 확인
        if shape_name in mapping:
            text_to_insert = mapping[shape_name]
            
            # 2) 텍스트를 넣을 수 있는 상자인지 확인
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear() # 기존 텍스트(제목을 입력하세요 등) 삭제
                
                p = text_frame.paragraphs[0]
                p.text = text_to_insert
                print(f"✅ [성공] '{shape_name}' 상자에 텍스트 입력됨.")
                matched_count += 1
            else:
                print(f"⚠️ [주의] '{shape_name}'은 이름은 맞지만 텍스트 상자가 아닙니다.")
        else:
            # 매칭 실패 시 (디버깅용)
            print(f"ℹ️ [Skip] PPT엔 '{shape_name}' 상자가 있는데, 보낼 데이터(JSON)엔 없습니다.")

    # 4. 결과 저장
    if matched_count > 0:
        prs.save(output_path)
        print(f"\n✨ 파일 생성 완료! '{output_path}'를 확인하세요.")
    else:
        print("\n❌ 경고: 매칭된 데이터가 하나도 없습니다. PPT 상자 이름(Key)을 확인하세요!")

# ====================================================
# 실행
# ====================================================
if __name__ == "__main__":
    create_ppt_file_test(mock_slide_data, TEMPLATE_PATH, OUTPUT_PATH)


# [수정된 버전] 3. Renderer 함수 (원본 레이아웃 이름 추적 방식)
def create_ppt_file(slide_data, template_path, output_path):
    prs = Presentation(template_path)
    
    # 1. LLM이 고른 레이아웃 가져오기
    # (Pydantic 객체 접근법 사용)
    try:
        target_index = slide_data.layout_index
        selected_layout = prs.slide_layouts[target_index]
    except (AttributeError, KeyError):
        # 딕셔너리로 들어올 경우 대비
        target_index = slide_data["layout_index"] if isinstance(slide_data, dict) else slide_data.layout_index
        selected_layout = prs.slide_layouts[target_index]
        
    # 2. 슬라이드 추가
    slide = prs.slides.add_slide(selected_layout)
    print(f"🎨 선택된 레이아웃: {selected_layout.name} (Index: {target_index})")

    # 데이터 매핑 준비
    if hasattr(slide_data, "content_mapping"):
        mapping = slide_data.content_mapping
    else:
        mapping = slide_data["content_mapping"]

    # 3. 데이터 매핑 (여기가 핵심 수정!!!)
    for shape in slide.placeholders:
        # 슬라이드 상자의 이름(shape.name)을 쓰는 게 아니라,
        # '번호(idx)'를 이용해서 '레이아웃의 원래 이름'을 찾아옵니다.
        try:
            shape_idx = shape.placeholder_format.idx
            # "설계도야, 이 번호(idx) 가진 상자 원래 이름이 뭐니?"
            original_name = selected_layout.placeholders[shape_idx].name
        except KeyError:
            # 혹시라도 못 찾으면 그냥 현재 이름 사용
            original_name = shape.name

        print(f"  🔍 확인 중: 슬라이드상 이름 '{shape.name}' -> 원본 이름 '{original_name}'")

        # 이제 '원본 이름'으로 매핑을 시도합니다.
        if original_name in mapping:
            content = mapping[original_name]
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content
                print(f"    ✅ 매칭 성공! 내용 입력 완료.")
        else:
            # 매칭 안 된 경우 (디버깅용 로그)
            pass 

    # 4. 저장
    prs.save(output_path)
    print(f"\n✨ 파일 생성 완료! {output_path}")


def create_ppt_file(slide_data, template_path, output_path):
    prs = Presentation(template_path)
    
    # 1. 레이아웃 가져오기
    try:
        target_index = slide_data.layout_index
        selected_layout = prs.slide_layouts[target_index]
    except (AttributeError, KeyError):
        target_index = slide_data["layout_index"] if isinstance(slide_data, dict) else slide_data.layout_index
        selected_layout = prs.slide_layouts[target_index]
        
    # 2. 슬라이드 추가
    slide = prs.slides.add_slide(selected_layout)
    print(f"🎨 선택된 레이아웃: {selected_layout.name} (Index: {target_index})")

    # 데이터 매핑 준비
    if hasattr(slide_data, "content_mapping"):
        mapping = slide_data.content_mapping
    else:
        mapping = slide_data["content_mapping"]

    # 3. 데이터 매핑 (안전한 Loop 방식 적용)
    for shape in slide.placeholders:
        # 슬라이드 상자의 번호표(idx) 확인
        shape_idx = shape.placeholder_format.idx
        
        # [핵심 수정] selected_layout.placeholders[shape_idx] 라고 쓰면 에러가 남!
        # 대신, 레이아웃의 상자들을 하나씩 돌면서 번호가 같은지 직접 확인합니다.
        
        original_name = shape.name # 못 찾을 경우를 대비한 기본값
        
        for layout_shape in selected_layout.placeholders:
            # 레이아웃 상자의 번호와 슬라이드 상자의 번호가 같으면?
            if layout_shape.placeholder_format.idx == shape_idx:
                original_name = layout_shape.name # 그 이름을 가져옴 (예: Body_Left)
                break
        
        print(f"  🔍 매핑 시도: 슬라이드(IDX:{shape_idx}) -> 원본이름 '{original_name}'")

        # 찾은 원본 이름으로 데이터 매핑
        if original_name in mapping:
            content = mapping[original_name]
            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content
                print(f"    ✅ 성공! 내용 입력됨.")
        else:
            # 매칭 안 된 경우
            pass 

    # 4. 저장
    prs.save(output_path)
    print(f"\n✨ 파일 생성 완료! {output_path}")

