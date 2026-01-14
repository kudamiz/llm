8from unstructured.partition.pptx import partition_pptx
import os

# 1. 경로 설정
output_image_dir = "./extracted_images"
os.makedirs(output_image_dir, exist_ok=True)

filename = "example_presentation.pptx"

# 2. PPT 파싱 (핵심 부분)
# extract_images_in_pdf=True 옵션은 PPTX에서도 작동하여 포함된 그림을 추출해줍니다.
elements = partition_pptx(
    filename=filename,
    extract_images_in_pdf=True,  # 이미지/차트 추출 활성화
    infer_table_structure=True,  # 표 구조(html) 추출 활성화
    image_output_dir_path=output_image_dir, # 추출된 이미지 저장 경로
)

# 3. 요소별 데이터 분류 (RAG용 데이터 전처리)
text_elements = []
table_elements = []
image_elements = []

for element in elements:
    # 요소의 타입 확인
    el_type = element.category
    
    if el_type == "Table":
        # 표는 HTML 메타데이터와 텍스트를 함께 저장
        table_elements.append({
            "text": element.text,
            "html": element.metadata.text_as_html,
            "page": element.metadata.page_number
        })
    
    elif el_type == "Image":
        # 이미지는 저장된 경로를 참조
        image_elements.append({
            "path": element.metadata.image_path,
            "page": element.metadata.page_number
        })
        
    elif el_type in ["Title", "NarrativeText", "ListItem"]:
        # 일반 텍스트
        text_elements.append({
            "text": element.text,
            "page": element.metadata.page_number
        })

print(f"텍스트 청크: {len(text_elements)}개")
print(f"추출된 표: {len(table_elements)}개")
print(f"추출된 이미지(차트 등): {len(image_elements)}개")


import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def iter_shapes(shapes):
    """그룹 안에 숨은 도형까지 샅샅이 뒤지는 재귀 함수"""
    for shape in shapes:
        # 1. 그룹인 경우: 재귀적으로 내부 진입
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape

def extract_images_from_pptx(pptx_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    prs = Presentation(pptx_path)
    image_count = 0

    print(f"이미지 추출 시작: {pptx_path}")

    for i, slide in enumerate(prs.slides):
        # 슬라이드 내의 모든 도형(그룹 포함)을 순회
        for shape in iter_shapes(slide.shapes):
            
            # 2. 그림(Picture)인 경우
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    # 이미지 바이너리 데이터 가져오기
                    image_bytes = image.blob
                    # 확장자 결정 (jpg, png 등)
                    ext = image.ext
                    
                    filename = f"slide_{i+1}_img_{image_count}.{ext}"
                    filepath = os.path.join(output_dir, filename)
                    
                    with open(filepath, "wb") as f:
                        f.write(image_bytes)
                        
                    print(f"  [저장됨] {filename}")
                    image_count += 1
                except Exception as e:
                    print(f"  [에러] 이미지 저장 실패: {e}")

    print(f"총 {image_count}개의 이미지를 추출했습니다.")

# --- 실행 ---
extract_images_from_pptx("example.pptx", "./extracted_images")

import os
import subprocess
from pdf2image import convert_from_path

def ppt_to_images_via_libreoffice(ppt_path, output_dir):
    """
    LibreOffice를 이용해 PPT를 이미지로 변환하는 함수
    1. PPT -> PDF (LibreOffice headless 모드 사용)
    2. PDF -> Images (pdf2image 사용)
    """
    
    # 0. 경로 설정
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    ppt_dir = os.path.dirname(ppt_path)
    ppt_filename = os.path.basename(ppt_path)
    pdf_filename = os.path.splitext(ppt_filename)[0] + ".pdf"
    pdf_path = os.path.join(output_dir, pdf_filename) # PDF도 일단 output_dir에 저장

    print(f"1. PDF 변환 시작: {ppt_filename}...")

    # 1. LibreOffice를 이용해 PDF로 변환 (터미널 명령어 실행)
    # --headless: 화면 없이 실행
    # --convert-to pdf: PDF로 변환
    # --outdir: 저장할 폴더
    command = [
        "libreoffice", 
        "--headless", 
        "--convert-to", "pdf", 
        "--outdir", output_dir, 
        ppt_path
    ]
    
    try:
        # subprocess로 리눅스 명령어 실행
        result = subprocess.run(command, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        # print(result.stdout.decode()) # 로그 보고 싶으면 주석 해제
    except subprocess.CalledProcessError as e:
        print(f"❌ LibreOffice 변환 실패: {e.stderr.decode()}")
        return []

    if not os.path.exists(pdf_path):
        print("❌ PDF 파일이 생성되지 않았습니다.")
        return []

    print("2. 이미지 변환 시작 (PDF -> Images)...")

    # 2. 변환된 PDF를 이미지로 쪼개기
    try:
        # dpi=300 : 고화질 설정 (OCR/VLM 인식률 높이려면 300 추천)
        images = convert_from_path(pdf_path, dpi=300)
        
        saved_image_paths = []
        for i, image in enumerate(images):
            # 슬라이드 번호는 1부터 시작
            image_filename = f"slide_{i+1}.jpg"
            save_path = os.path.join(output_dir, image_filename)
            
            image.save(save_path, "JPEG")
            saved_image_paths.append(save_path)
            print(f"  - 저장됨: {save_path}")
            
        print(f"✅ 변환 완료! 총 {len(saved_image_paths)}장")
        
        # (선택) 중간에 만든 PDF는 삭제하고 싶다면:
        # os.remove(pdf_path)
        
        return saved_image_paths

    except Exception as e:
        print(f"❌ 이미지 변환 실패: {e}")
        return []

# --- 실행 예시 ---
# ppt_file = "./data/defect_report.pptx"
# output_folder = "./extracted_images/report_01"

# images = ppt_to_images_via_libreoffice(ppt_file, output_folder)


import os
import subprocess
import base64
import json
from pdf2image import convert_from_path
from unstructured.partition.pptx import partition_pptx
import nltk

# =============================================================================
# 0. 환경 설정 (NLTK 오프라인 경로 & VLM 클라이언트)
# =============================================================================

# NLTK 데이터 경로 강제 지정 (서버 오프라인 이슈 해결용)
nltk_data_path = os.path.abspath("./nltk_data")
if nltk_data_path not in nltk.data.path:
    nltk.data.path.insert(0, nltk_data_path)

# (예시) VLM 호출 함수 - 실제 사용하는 모델(GPT-4o, Gemini) API로 교체 필요
def call_vlm_api(image_path, raw_text_hint):
    """
    이미지와 힌트 텍스트를 받아 시각적 분석 결과를 반환하는 가상의 함수
    """
    # 실제 구현 시: openai.ChatCompletion.create(...) 또는 langchain 등 사용
    # 프롬프트: "텍스트 읽지 말고, 불량 위치/형태/그래프 추이 등 시각적 맥락만 설명해."
    
    # --- Dummy Return (테스트용) ---
    return f"(VLM 분석 결과) 이 이미지는 오른쪽 상단에 붉은 원으로 표시된 크랙을 보여줌. 텍스트 힌트 '{raw_text_hint[:10]}...'와 관련 있어 보임."

# =============================================================================
# 1. PPT -> 이미지 변환 (LibreOffice + pdf2image)
# =============================================================================
def convert_ppt_to_images(ppt_path, output_dir):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    ppt_filename = os.path.basename(ppt_path)
    pdf_filename = os.path.splitext(ppt_filename)[0] + ".pdf"
    
    print(f"🔄 [1/4] 이미지 변환 시작: {ppt_filename}")
    
    # 1-1. LibreOffice로 PDF 변환
    cmd = ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, ppt_path]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    pdf_path = os.path.join(output_dir, pdf_filename)
    if not os.path.exists(pdf_path):
        print("❌ PDF 변환 실패")
        return {}

    # 1-2. PDF -> 이미지 리스트 변환
    images = convert_from_path(pdf_path, dpi=300) # 고화질
    image_map = {} # {page_num: image_path}
    
    for i, img in enumerate(images):
        page_num = i + 1
        img_name = f"slide_{page_num}.jpg"
        save_path = os.path.join(output_dir, img_name)
        img.save(save_path, "JPEG")
        image_map[page_num] = save_path
        
    print(f"✅ 총 {len(image_map)}장 이미지 변환 완료")
    return image_map

# =============================================================================
# 2. 텍스트 & 표 추출 (Unstructured)
# =============================================================================
def extract_text_data(ppt_path):
    print(f"mining [2/4] 텍스트 및 표 추출 시작...")
    
    # 이미지 추출은 LibreOffice로 하므로 여기선 텍스트만 빠르게 추출
    elements = partition_pptx(
        filename=ppt_path,
        include_page_breaks=False,
        infer_table_structure=True # 표 구조 분석 켬
    )
    
    slides_data = {} # {page_num: {"text": "", "tables": []}}
    
    for el in elements:
        page_num = el.metadata.page_number
        if page_num not in slides_data:
            slides_data[page_num] = {"text": [], "tables": []}
            
        if el.category == "Table":
            # 표는 HTML 형태로 저장 + 텍스트에도 추가
            slides_data[page_num]["tables"].append(el.metadata.text_as_html)
            slides_data[page_num]["text"].append(el.text)
        elif el.category in ["Title", "NarrativeText", "ListItem"]:
            slides_data[page_num]["text"].append(el.text)
            
    # 리스트를 하나의 문자열로 합치기
    for page in slides_data:
        slides_data[page]["text"] = "\n".join(slides_data[page]["text"])
        
    print(f"✅ {len(slides_data)}페이지 텍스트 추출 완료")
    return slides_data

# =============================================================================
# 3. 데이터 병합 및 구조화 (OpenSearch Schema + Context Injection)
# =============================================================================
def build_rag_documents(ppt_path, image_map, text_data):
    print(f"🧩 [3/4] VLM 분석 및 데이터 구조화 (Context Injection)...")
    
    filename = os.path.basename(ppt_path)
    final_docs = []
    
    # 글로벌 맥락 (문서 전체 주제 - 실제론 LLM으로 파일 전체 요약 추천)
    global_context = f"문서: {filename}, 주제: 반도체 불량 분석 리포트"
    
    # 이전 슬라이드 요약 (Context Flow) - 초기값
    prev_slide_summary = "첫 페이지입니다."
    
    # 페이지 순서대로 처리
    all_pages = sorted(list(set(image_map.keys()) | set(text_data.keys())))
    
    for page in all_pages:
        # 데이터 가져오기 (없으면 빈값 처리)
        raw_text = text_data.get(page, {}).get("text", "")
        img_path = image_map.get(page)
        
        # --- [Step 3. VLM 분석] ---
        # 텍스트 힌트를 주어 VLM이 이미지를 더 잘 보게 함
        if img_path:
            vlm_desc = call_vlm_api(img_path, raw_text_hint=raw_text)
        else:
            vlm_desc = "이미지 없음"
            
        # --- [Step 4. OpenSearch용 데이터 조립] ---
        
        # 4-1. 검색용 텍스트 (Vector Embedding 대상)
        # 글로벌 주제 + 이전 장 내용 + 현재 이미지 설명 + 현재 텍스트
        search_context_blob = f"""
        [GLOBAL TOPIC] {global_context}
        [PREV CONTEXT] {prev_slide_summary}
        [VISUAL DESC] {vlm_desc}
        [RAW CONTENT] {raw_text}
        """
        
        # 4-2. 최종 스키마 (Flat JSON)
        doc = {
            "id": f"{filename}_p{page}",           # ID
            "filename": filename,                  # 필터용
            "page": page,                          # 정렬용
            "search_context": search_context_blob, # ★ 임베딩할 핵심 데이터
            "display_content": raw_text,           # 화면 표시용 텍스트
            "image_path": img_path,                # 화면 표시용 이미지 경로
            "vlm_analysis": vlm_desc               # 화면 표시용 분석글
        }
        
        final_docs.append(doc)
        
        # 4-3. 다음 루프를 위해 현재 내용을 요약하여 '이전 장 정보'로 업데이트
        # (간단히 앞부분만 잘라서 씀, 실제론 LLM 요약 추천)
        prev_slide_summary = (raw_text + vlm_desc)[:200].replace("\n", " ")
        
    print(f"🎉 [4/4] 최종 데이터 생성 완료: {len(final_docs)}개 문서")
    return final_docs

# =============================================================================
# 메인 실행
# =============================================================================
if __name__ == "__main__":
    # 설정
    target_ppt = "./data/sample_defect.pptx"
    image_out_dir = "./extracted_images"
    
    # 1. 이미지 변환
    images = convert_ppt_to_images(target_ppt, image_out_dir)
    
    # 2. 텍스트 추출
    texts = extract_text_data(target_ppt)
    
    # 3. RAG용 데이터 생성 (VLM 포함)
    rag_ready_data = build_rag_documents(target_ppt, images, texts)
    
    # 결과 확인 (첫 번째 슬라이드만)
    if rag_ready_data:
        print("\n--- [결과 미리보기 (첫장)] ---")
        print(json.dumps(rag_ready_data[0], indent=2, ensure_ascii=False))
        
    # TODO: 여기서 rag_ready_data를 loop 돌며 Embedding API 호출 -> OpenSearch Bulk Insert
