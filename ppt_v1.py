import os
from typing import TypedDict, Dict, List
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from langchain_openai import ChatOpenAI
from pydantic import BaseModel, Field
from langgraph.graph import StateGraph, END

# ==============================================================================
# 1. Helper Functions (핵심 엔진)
# ==============================================================================

def get_real_placeholder_name(shape):
    """Placeholder의 진짜 이름(Master Layout상의 이름) 추적"""
    try:
        if not shape.is_placeholder: return shape.name
        layout = shape.part.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == shape.placeholder_format.idx:
                return ph.name
        return shape.name
    except: return shape.name

def replace_text_preserving_style(shape, new_text):
    """서식 보존 텍스트 교체 (빈 칸 방어 로직 포함)"""
    if not shape.has_text_frame: return
    tf = shape.text_frame
    new_text_str = str(new_text)

    # 빈 칸이면 그냥 넣기
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        tf.text = new_text_str
        return

    # 스타일 백업
    p = tf.paragraphs[0]
    sample_run = p.runs[0]
    font_name = sample_run.font.name
    font_size = sample_run.font.size
    font_bold = sample_run.font.bold
    font_color_rgb = sample_run.font.color.rgb if hasattr(sample_run.font.color, 'rgb') else None

    # 교체
    p.clear()
    new_run = p.add_run()
    new_run.text = new_text_str
    
    # 복원
    if font_name: new_run.font.name = font_name
    if font_size: new_run.font.size = font_size
    if font_bold is not None: new_run.font.bold = font_bold
    if font_color_rgb: new_run.font.color.rgb = font_color_rgb

def smart_fill_placeholders(slide, data_dict):
    """이름 매칭 -> 타입 매칭 순으로 빈칸 채우기"""
    norm_data = {k.lower(): v for k, v in data_dict.items()}
    
    for shape in slide.placeholders:
        real_name = get_real_placeholder_name(shape).lower()
        ph_type = shape.placeholder_format.type
        target = None
        
        print(f"   🔍 슬라이드 칸 분석: {real_name} (Type: {ph_type})")

        # 1. 이름 매칭 (Exact & Partial)
        # 데이터 키가 슬라이드 칸 이름에 포함되면 매칭 (예: 'title' in 'master title')
        for k, v in norm_data.items():
            if k in real_name:
                target = v
                break
        
        # 2. 타입 매칭 (Fallback)
        if not target:
            if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                target = norm_data.get("title") or norm_data.get("subject")
            elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                target = norm_data.get("subtitle")
            elif ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                target = norm_data.get("content") or norm_data.get("body")
            elif ph_type == PP_PLACEHOLDER.DATE:
                target = norm_data.get("date")

        if target:
            replace_text_preserving_style(shape, target)
            print(f"      ✅ 채우기 성공: '{target}'")

# ==============================================================================
# 2. State & Schema
# ==============================================================================

class StaticState(TypedDict):
    # Input
    template_path: str
    output_path: str
    target_layout_name: str  # 사용자가 고른 레이아웃 (예: "Title_Slide")
    user_instruction: str    # 사용자 요청 (예: "제목은 실적보고, 부제는 1팀")
    
    # Internal
    layout_info: str         # Scanner가 찾은 Placeholder 목록
    generated_data: dict     # Content가 만든 데이터 {"Title": "...", "Subtitle": "..."}

class ContentData(BaseModel):
    # 어떤 키가 들어올지 모르므로 유연한 Dict 형태로 정의
    fields: Dict[str, str] = Field(
        ..., 
        description="Placeholder 이름과 매핑될 데이터. 예: {'Title': '제목', 'Content': '내용'}"
    )

# ==============================================================================
# 3. Nodes (Scanner -> Content -> Renderer)
# ==============================================================================

def scanner_node(state: StaticState):
    print("\n--- [1] Scanner: 레이아웃 분석 ---")
    prs = Presentation(state["template_path"])
    target_name = state["target_layout_name"]
    
    found_layout = None
    for layout in prs.slide_layouts:
        if layout.name == target_name:
            found_layout = layout
            break
            
    if not found_layout:
        raise ValueError(f"❌ 템플릿에서 '{target_name}' 레이아웃을 찾을 수 없습니다.")
        
    # Placeholder 이름 수집
    ph_names = [get_real_placeholder_name(ph) for ph in found_layout.placeholders]
    
    info = f"""
    Layout Name: {target_name}
    Available Placeholders: {', '.join(ph_names)}
    """
    print(info)
    
    return {"layout_info": info}

def content_node(state: StaticState):
    print("\n--- [2] Content: 내용 매핑 ---")
    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    structured_llm = llm.with_structured_output(ContentData)
    
    system_prompt = f"""
    당신은 PPT 슬라이드 작성기입니다.
    [사용자 요청]을 분석하여 [가능한 칸]에 넣을 데이터를 JSON으로 만드세요.
    
    [가능한 칸 (Placeholders)]
    {state['layout_info']}
    
    [사용자 요청]
    {state['user_instruction']}
    
    [작성 규칙]
    - 'fields' 딕셔너리에 "Placeholder이름": "내용" 형식으로 담으세요.
    - 가능한 칸의 이름과 비슷하게 Key를 잡으면 매핑이 잘 됩니다. (예: Title -> Title)
    """
    
    res = structured_llm.invoke(system_prompt)
    return {"generated_data": res.fields}

def renderer_node(state: StaticState):
    print("\n--- [3] Renderer: 파일 생성 ---")
    prs = Presentation(state["template_path"])
    target_name = state["target_layout_name"]
    
    # 해당 레이아웃으로 슬라이드 1장 추가
    layout = next(l for l in prs.slide_layouts if l.name == target_name)
    slide = prs.slides.add_slide(layout)
    
    # 데이터 채우기
    data = state["generated_data"]
    smart_fill_placeholders(slide, data)
    
    prs.save(state["output_path"])
    print(f"🎉 생성 완료: {state['output_path']}")
    return {"output_path": state["output_path"]}

# ==============================================================================
# 4. Graph Wiring
# ==============================================================================

workflow = StateGraph(StaticState)

workflow.add_node("scanner", scanner_node)
workflow.add_node("content", content_node)
workflow.add_node("renderer", renderer_node)

workflow.set_entry_point("scanner")
workflow.add_edge("scanner", "content")
workflow.add_edge("content", "renderer")
workflow.add_edge("renderer", END)

app = workflow.compile()

# ==============================================================================
# 5. Test Execution
# ==============================================================================

if __name__ == "__main__":
    # 사용 예시
    inputs = {
        "template_path": "template.pptx",  # 템플릿 파일 경로
        "output_path": "static_result.pptx",
        
        # [사용자 입력 1] 어떤 슬라이드를 쓸 건지?
        "target_layout_name": "Title_Slide", 
        
        # [사용자 입력 2] 무슨 내용을 넣을 건지?
        "user_instruction": "제목은 '2024년 4분기 실적보고'로 하고, 부제목은 '영업1팀 김철수'로 해줘. 날짜는 2024-12-01."
    }
    
    try:
        app.invoke(inputs)
    except Exception as e:
        print(f"에러 발생: {e}")
