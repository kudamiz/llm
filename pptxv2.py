import os
import io
from typing import TypedDict, Optional, Dict, List
from pydantic import BaseModel, Field

# 라이브러리 임포트
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langgraph.graph import StateGraph, END

# ====================================================
# [설정 1] 레이아웃별 상세 가이드 (Router & Generator용)
# PPT 슬라이드 마스터의 '레이아웃 이름'과 Key가 일치해야 합니다.
# ====================================================
LAYOUT_DETAILS = {
    "Comparison_Slide": {
        "description": "두 가지 대상을 비교할 때 사용합니다. (예: 경쟁사 비교, 전후 비교)",
        "fields": {
            "Title": "비교 주제를 명확히 작성",
            "Left_Item": "A대상의 장점을 3줄 요약 (개조식)",
            "Right_Item": "B대상의 장점을 3줄 요약 (개조식)",
            "table_spec": "반드시 정량적 수치(%)가 포함된 표 데이터 작성"
        }
    },
    "Project_Overview": {
        "description": "프로젝트의 개요나 성과를 보고할 때 사용합니다.",
        "fields": {
            "Title": "임팩트 있는 헤드라인",
            "Goal": "핵심 목표 3가지를 개조식으로 작성",
            "Effect": "기대 효과를 구체적인 수치로 표현",
            "image_main": "프로젝트 관련 고화질 이미지 사용"
        }
    },
    # 설정이 없는 레이아웃은 기본값(설명 없음)으로 처리됩니다.
}

# ====================================================
# [설정 2] 헬퍼 함수 (PPT 조작 도구들)
# ====================================================
def insert_multiline_text(shape, content):
    """줄바꿈 문자(\\n)를 인식하여 단락을 나누어 입력"""
    if not shape.has_text_frame: return
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.TOP # 상단 정렬

    lines = content.strip().split('\n')
    if lines:
        p = text_frame.paragraphs[0]
        p.text = lines[0]
        # (Mac용 폰트 설정)
        p.font.name = 'AppleSDGothicNeo' 
        p.font.size = Pt(18)

    for line in lines[1:]:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.name = 'AppleSDGothicNeo'
        p.font.size = Pt(18)

def parse_table_string(text_data):
    """ 'A|B\\nC|D' 문자열을 리스트로 변환 """
    rows = []
    for line in text_data.strip().split('\n'):
        if "|" in line:
            cols = [c.strip() for c in line.split('|')]
            if set(cols[0]) <= {'-', ' '}: continue
            rows.append(cols)
    return rows

def insert_styled_table(shape, content_string):
    """표 생성 및 스타일/폰트 적용"""
    table_data = parse_table_string(content_string)
    rows = len(table_data)
    cols = len(table_data[0]) if rows > 0 else 0
    if rows == 0: return

    try:
        graphic_frame = shape.insert_table(rows=rows, cols=cols)
        table = graphic_frame.table
        # PPT 기본 스타일 ID 적용 (테마 색상 추종)
        table.table_style_id = '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'

        for r in range(rows):
            for c in range(cols):
                cell = table.cell(r, c)
                cell.text = str(table_data[r][c])
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = 'AppleSDGothicNeo' # Mac 호환 폰트
                    p.alignment = PP_ALIGN.CENTER
                    if r == 0: p.font.bold = True
    except AttributeError:
        print(f"      ❌ 에러: '{shape.name}'은 표 타입이 아닙니다.")

# ====================================================
# [LangGraph] 1. State 정의
# ====================================================
class AgentState(TypedDict):
    # Input
    user_query: str
    template_path: str
    output_path: str
    image_files: Dict[str, bytes] # 파일명: 바이너리 데이터
    
    # Internal State
    template_guide: Optional[str]
    slide_data: Optional[object] # Pydantic Object
    final_message: Optional[str]

# ====================================================
# [LangGraph] 2. Node 정의
# ====================================================

def scanner_node(state: AgentState):
    """[Step 1] 템플릿 분석 및 가이드(Router+Rules) 생성"""
    print("\n--- [Node 1] Scanner: 템플릿 분석 중 ---")
    pptx_path = state["template_path"]
    
    if not os.path.exists(pptx_path):
        return {"template_guide": "ERROR: 템플릿 파일이 없습니다."}

    prs = Presentation(pptx_path)
    guide_lines = ["=== 템플릿 선택 및 작성 가이드 ==="]

    for i, layout in enumerate(prs.slide_layouts):
        layout_name = layout.name
        
        # 설정 가져오기
        details = LAYOUT_DETAILS.get(layout_name, {})
        desc = details.get("description", "용도 설명 없음 (이름 참고)")
        field_rules = details.get("fields", {})

        # 헤더 작성 (Router 역할)
        guide_lines.append(f"\n[Layout Index: {i}] 이름: {layout_name}\n   💡 용도: {desc}")

        # 칸별 규칙 작성 (Generator 역할)
        for shape in layout.placeholders:
            p_name = shape.name
            rule = field_rules.get(p_name, "")
            if rule:
                info = f"   - 칸 '{p_name}': ⭐[규칙: {rule}]"
            else:
                info = f"   - 칸 '{p_name}': (자유 작성)"
            guide_lines.append(info)
            
    return {"template_guide": "\n".join(guide_lines)}

def planner_node(state: AgentState):
    """[Step 2] AI 기획 (Layout 선택 + Content 생성)"""
    print("--- [Node 2] Planner: AI 기획 중 ---")
    guide = state["template_guide"]
    
    if "ERROR" in guide:
        print("❌ 템플릿 오류로 중단")
        return {"slide_data": None}

    # 출력 포맷 정의
    class SlideOutput(BaseModel):
        layout_index: int = Field(..., description="선택한 레이아웃 인덱스")
        content_mapping: Dict[str, str] = Field(..., description="Placeholder 이름과 내용 매핑")
        reason: str = Field(..., description="레이아웃 선택 이유")

    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    structured_llm = llm.with_structured_output(SlideOutput)
    
    system_prompt = """
    당신은 PPT 생성 전문가입니다. [템플릿 가이드]를 분석하여 사용자 요청에 가장 적합한 레이아웃을 하나 선택하고 내용을 채우세요.
    
    [작성 원칙]
    1. '💡 용도'를 보고 가장 적절한 레이아웃을 선택(Router)하세요.
    2. '⭐[규칙:...]'이 있는 칸은 반드시 해당 규칙을 지켜서 내용을 작성하세요.
    3. 텍스트 줄바꿈이 필요하면 '\\n'을 사용하세요.
    4. 이미지는 업로드된 파일명을 값으로 넣으세요. (예: sample.jpg)
    5. 표는 '헤더|헤더\\n값|값' 형태의 문자열로 작성하세요.

    [템플릿 가이드]
    {guide}
    """
    
    chain = ChatPromptTemplate.from_messages([("system", system_prompt), ("human", "{input}")]) | structured_llm
    result = chain.invoke({"guide": guide, "input": state["user_query"]})
    
    print(f"👉 선택된 레이아웃: {result.layout_index}번 (이유: {result.reason})")
    return {"slide_data": result}

def renderer_node(state: AgentState):
    """[Step 3] 파일 생성 (이미지 Binary 삽입 포함)"""
    print("--- [Node 3] Renderer: PPT 생성 중 ---")
    data = state["slide_data"]
    if not data: return {"final_message": "데이터 없음"}

    prs = Presentation(state["template_path"])
    slide = prs.slides.add_slide(prs.slide_layouts[data.layout_index])
    mapping = data.content_mapping
    
    # 매핑 로직
    for shape in slide.placeholders:
        shape_idx = shape.placeholder_format.idx
        # 원본 이름 찾기
        original_name = shape.name
        for layout_shape in prs.slide_layouts[data.layout_index].placeholders:
            if layout_shape.placeholder_format.idx == shape_idx:
                original_name = layout_shape.name
                break
        
        if original_name in mapping:
            content = mapping[original_name]
            
            # [이미지] 메모리 내 Binary 처리
            if original_name.lower().startswith("image_"):
                # 파일명으로 Binary 데이터 찾기
                image_bytes = state["image_files"].get(content)
                if image_bytes:
                    try:
                        # BytesIO로 변환하여 삽입
                        image_stream = io.BytesIO(image_bytes)
                        shape.insert_picture(image_stream)
                        print(f"    🖼️ 이미지 삽입 성공: {content}")
                    except Exception as e:
                        print(f"    ❌ 이미지 처리 실패: {e}")
                else:
                    print(f"    ⚠️ 경고: '{content}' 파일이 image_files 목록에 없습니다.")

            # [표] 스타일 적용
            elif original_name.lower().startswith("table_"):
                insert_styled_table(shape, content)
                print("    📊 표 삽입 완료")
                
            # [텍스트] 줄바꿈 지원
            else:
                insert_multiline_text(shape, content)

    prs.save(state["output_path"])
    msg = f"완료! 저장 경로: {state['output_path']}"
    print(f"✨ {msg}")
    return {"final_message": msg}

# ====================================================
# [LangGraph] 3. 그래프 구성
# ====================================================
workflow = StateGraph(AgentState)
workflow.add_node("scanner", scanner_node)
workflow.add_node("planner", planner_node)
workflow.add_node("renderer", renderer_node)

workflow.set_entry_point("scanner")
workflow.add_edge("scanner", "planner")
workflow.add_edge("planner", "renderer")
workflow.add_edge("renderer", END)

app = workflow.compile()

# ====================================================
# [실행] 테스트 시뮬레이션
# ====================================================
if __name__ == "__main__":
    # [준비물 1] 템플릿 파일 체크
    template_file = "my_template.pptx"
    if not os.path.exists(template_file):
        print("❌ 테스트를 위해 'my_template.pptx' 파일이 필요합니다.")
        exit()

    # [준비물 2] 이미지 파일 Binary 로딩 (API 업로드 흉내)
    image_filename = "sample_image.jpg"
    image_data = {}
    
    if os.path.exists(image_filename):
        with open(image_filename, "rb") as f:
            image_data[image_filename] = f.read() # Bytes 형태로 저장
            print(f"📂 이미지 로드 완료: {image_filename}")
    else:
        print(f"⚠️ 경고: '{image_filename}' 없음. 이미지 기능 테스트 불가.")

    # [입력] 사용자 요청
    inputs = {
        "user_query": f"이번 프로젝트 성과를 보고하려고 해. 목표 달성과 기대효과를 강조해주고, 메인 사진으로 '{image_filename}'을 넣어줘.",
        "template_path": template_file,
        "output_path": "final_result.pptx",
        "image_files": image_data # Dict[파일명, Bytes]
    }
    
    # [실행]
    print("🚀 에이전트 시작...")
    for output in app.stream(inputs):
        pass
