import io
import os
from typing import List, Dict, Any, Literal, Optional, TypedDict, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pydantic import BaseModel, Field
from langchain_openai import ChatOpenAI
from langgraph.graph import StateGraph, END

# ====================================================
# [Config] 템플릿 레지스트리 (Hybrid 규칙 정의)
# ====================================================
# Scanner가 1차적으로 참고하되, 실제 파일의 Anchor도 함께 읽습니다.
TEMPLATE_REGISTRY = {
    "Title_Slide": {
        "type": "static",
        "desc": "표지 슬라이드",
        "rules": {"Title": "20자 이내, 임팩트 있게", "Subtitle": "날짜/발표자 포함"}
    },
    "Content_List": {
        "type": "static",
        "desc": "목차 및 아젠다",
        "rules": {"Content": "개조식으로 요약"}
    },
    "Dynamic_Split": {
        "type": "dynamic",
        "desc": "좌우 비교 및 분석 (차트/텍스트 혼합)",
        "rules": {"Guide_Left": "차트 배치 추천", "Guide_Right": "핵심 요약 텍스트"}
    },
    "Dynamic_Full": {
        "type": "dynamic",
        "desc": "대형 데이터 시각화",
        "rules": {"Guide_Main": "복잡한 표나 큰 차트"}
    }
}

# ====================================================
# [Helper 1] 서식 보존 텍스트 교체 (Run-Level) - NEW!
# ====================================================
def fill_placeholder_preserving_style(shape, new_text):
    """
    기존 텍스트의 폰트/색상/크기를 최대한 유지하며 내용을 교체합니다.
    첫 번째 문단의 첫 번째 Run 스타일을 복사하여 적용합니다.
    """
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    if not text_frame.paragraphs:
        text_frame.text = new_text # 문단 없으면 그냥 넣음
        return

    # 첫 번째 문단의 첫 번째 런(Run) 스타일 가져오기
    p = text_frame.paragraphs[0]
    if p.runs:
        r = p.runs[0]
        font_name = r.font.name
        font_size = r.font.size
        font_bold = r.font.bold
        font_color = r.font.color.rgb if hasattr(r.font.color, 'rgb') else None
    else:
        # 런이 없으면 그냥 텍스트 교체
        text_frame.text = new_text
        return

    # 텍스트 교체 (기존 내용 싹 지우고 새로 씀)
    text_frame.clear() 
    new_p = text_frame.paragraphs[0]
    new_run = new_p.add_run()
    new_run.text = new_text

    # 스타일 복원
    if font_name: new_run.font.name = font_name
    if font_size: new_run.font.size = font_size
    if font_bold is not None: new_run.font.bold = font_bold
    if font_color: new_run.font.color.rgb = font_color

# ====================================================
# [Helper 2] Placeholder 이름 역추적
# ====================================================
def get_real_ph_name(shape):
    try:
        idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph.name
        return shape.name
    except:
        return shape.name

# ====================================================
# [Helper 3] Dynamic Drawing Tools
# ====================================================
def draw_chart(slide, x, y, w, h, data):
    chart_data = CategoryChartData()
    chart_data.categories = data.get('labels', [])
    chart_data.add_series('Series 1', data.get('values', []))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    ).chart
    if data.get('title'):
        chart.chart_title.text_frame.text = data['title']

def draw_table(slide, x, y, w, h, rows):
    if not rows: return
    r_cnt, c_cnt = len(rows), len(rows[0])
    table = slide.shapes.add_table(r_cnt, c_cnt, x, y, w, h).table
    for r in range(r_cnt):
        for c in range(c_cnt):
            table.cell(r, c).text = str(rows[r][c])

def draw_text_box(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.text = text
    tb.text_frame.word_wrap = True




# [State] 그래프 전체에서 공유할 메모리
class AgentState(TypedDict):
    user_query: str           # 사용자 요청
    template_path: str        # PPTX 경로
    output_path: str          # 저장 경로
    
    template_summary: str     # Node 2용 (간략 가이드)
    template_details: str     # Node 3용 (상세 규칙)
    
    skeleton_plan: List[dict] # Node 2 결과 (뼈대)
    slide_data: List[dict]    # Node 3 결과 (최종 데이터)
    
    # Reviewer Loop용
    review_status: str        # PASS / FAIL
    review_feedback: str      # 피드백 내용
    retry_count: int          # 재시도 횟수

# [Schema 1] Structure Node용
class SlideSkeleton(BaseModel):
    layout_index: int
    slide_type: Literal["static", "dynamic"]
    topic: str

class Storyboard(BaseModel):
    plan: List[SlideSkeleton]

# [Schema 2] Content Node용 (유니버설 데이터 모델)
class ComponentData(BaseModel):
    text_content: Optional[str] = None
    table_rows: Optional[List[List[str]]] = None
    chart_labels: Optional[List[str]] = None
    chart_values: Optional[List[float]] = None
    chart_title: Optional[str] = None

class SlideComponent(BaseModel):
    type: Literal["text", "table", "chart", "image"]
    position: str
    data: ComponentData

class SlideContent(BaseModel):
    type: Literal["static", "dynamic"]
    layout_index: int
    # Static & Dynamic 공통 (제목 등)
    common_fields: Dict[str, str] = Field(default_factory=dict)
    # Dynamic 전용
    components: List[SlideComponent] = Field(default_factory=list)

class PresentationPlan(BaseModel):
    slides: List[SlideContent]

# [Schema 3] Reviewer Node용
class ReviewResult(BaseModel):
    status: Literal["PASS", "FAIL"]
    feedback: str


def scanner_node(state: AgentState):
    prs = Presentation(state["template_path"])
    summary_lines = []
    detail_lines = []

    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name
        
        # 레지스트리에 없으면 기본값 처리 (Hybrid)
        config = TEMPLATE_REGISTRY.get(name, {"type": "static", "desc": "일반 레이아웃", "rules": {}})
        
        # 1. Summary (Structure용)
        summary_lines.append(f"[Index {i}] {name} ({config['type']}) : {config['desc']}")
        
        # 2. Details (Content용)
        info = f"\n[Layout {i}] {name} ({config['type']})"
        rules = config.get("rules", {})
        
        # (A) Static Placeholders
        ph_names = [get_real_ph_name(ph) for ph in layout.placeholders]
        if ph_names:
            info += f"\n   - 입력칸: {', '.join(ph_names)}"
            
        # (B) Dynamic Anchors (실제 파일 조회)
        anchors = [s.name for s in layout.shapes if s.name.startswith("Guide_")]
        if anchors:
            info += f"\n   - 앵커: {', '.join(anchors)}"
            
        # (C) 규칙 매핑
        info += "\n   - 작성 규칙:"
        for key, rule in rules.items():
            info += f"\n     * {key}: {rule}"
            
        detail_lines.append(info)

    return {
        "template_summary": "\n".join(summary_lines),
        "template_details": "\n".join(detail_lines),
        "retry_count": 0, # 초기화
        "review_feedback": ""
    }


def structure_node(state: AgentState):
    print("--- [Node 2] Structure: 스토리보드 기획 ---")
    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    structured_llm = llm.with_structured_output(Storyboard)
    
    prompt = f"""
    사용자 요청: {state['user_query']}
    
    [템플릿 목록]
    {state['template_summary']}
    
    위 템플릿을 활용해 논리적인 슬라이드 목차를 기획하세요.
    """
    res = structured_llm.invoke(prompt)
    return {"skeleton_plan": [s.model_dump() for s in res.plan]}



def content_node(state: AgentState):
    print(f"--- [Node 3] Content: 내용 작성 (Retry: {state['retry_count']}) ---")
    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    structured_llm = llm.with_structured_output(PresentationPlan)
    
    skeletons = state["skeleton_plan"]
    details = state["template_details"]
    feedback = state["review_feedback"]
    
    system_prompt = f"""
    당신은 PPT 콘텐츠 작가입니다.
    기획안에 따라 각 슬라이드의 데이터를 작성하세요.
    
    [기획안]
    {skeletons}
    
    [템플릿 상세 규칙 (준수 필수)]
    {details}
    """
    
    if feedback and feedback != "Good":
        system_prompt += f"\n\n🚨 [수정 요청] 이전 작성 내용에 문제가 있습니다:\n{feedback}\n이 지적사항을 반영해 처음부터 다시 작성하세요."

    res = structured_llm.invoke(system_prompt)
    
    # Pydantic -> Dict 변환
    return {"slide_data": [s.model_dump() for s in res.slides]}


def reviewer_node(state: AgentState):
    print("--- [Node 4] Reviewer: 품질 검수 ---")
    
    # 3회 이상 실패 시 강제 통과
    if state["retry_count"] >= 3:
        print("   ⚠️ 재시도 횟수 초과 -> 강제 PASS")
        return {"review_status": "PASS", "review_feedback": "Max retries"}

    llm = ChatOpenAI(model="gpt-4o", temperature=0)
    structured_llm = llm.with_structured_output(ReviewResult)
    
    prompt = f"""
    [검수 기준]
    {state['template_details']}
    
    [작성된 데이터]
    {state['slide_data']}
    
    위 데이터가 규칙을 준수했는지 검사하세요.
    - 글자 수 제한, 필수 데이터(labels, values) 누락 여부 확인.
    - 문제가 있으면 FAIL과 피드백을, 없으면 PASS를 반환하세요.
    """
    
    res = structured_llm.invoke(prompt)
    print(f"   ⚖️ 판정: {res.status}")
    
    return {
        "review_status": res.status,
        "review_feedback": res.feedback,
        "retry_count": state["retry_count"] + 1
    }


def renderer_node(state: AgentState):
    print("--- [Node 5] Renderer: 파일 생성 ---")
    prs = Presentation(state["template_path"])
    
    for plan in state["slide_data"]:
        layout_idx = plan["layout_index"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        
        # [A] Static & Common Fields (서식 보존 교체)
        common = plan.get("common_fields", {})
        for shape in slide.placeholders:
            real_name = get_real_ph_name(shape)
            if shape.placeholder_format.type == 1: real_name = "Title" # 제목 강제 매핑
            
            if real_name in common:
                # NEW: 스타일 유지하며 교체 함수 사용
                fill_placeholder_preserving_style(shape, common[real_name])
                
        # [B] Dynamic Components
        if plan["type"] == "dynamic":
            layout = prs.slide_layouts[layout_idx]
            # 앵커 찾기 (Layout에서 조회)
            anchors = {s.name: (s.left, s.top, s.width, s.height) 
                       for s in layout.shapes if s.name.startswith("Guide_")}
            
            for comp in plan.get("components", []):
                pos = comp["position"]
                data = comp["data"]
                
                if pos in anchors:
                    x, y, w, h = anchors[pos]
                    c_type = comp["type"]
                    
                    if c_type == "text":
                        draw_text_box(slide, x, y, w, h, data["text_content"])
                    elif c_type == "table":
                        draw_table(slide, x, y, w, h, data["table_rows"])
                    elif c_type == "chart":
                        chart_d = {
                            "labels": data["chart_labels"],
                            "values": data["chart_values"],
                            "title": data["chart_title"]
                        }
                        draw_chart(slide, x, y, w, h, chart_d)
                        
    prs.save(state["output_path"])
    print(f"🎉 생성 완료: {state['output_path']}")
    return {"output_path": state["output_path"]}


def route_after_review(state: AgentState):
    if state["review_status"] == "FAIL":
        return "content" # 재작성
    return "renderer"    # 통과

workflow = StateGraph(AgentState)

# 노드 등록
workflow.add_node("scanner", scanner_node)
workflow.add_node("structure", structure_node)
workflow.add_node("content", content_node)
workflow.add_node("reviewer", reviewer_node)
workflow.add_node("renderer", renderer_node)

# 흐름 연결
workflow.set_entry_point("scanner")
workflow.add_edge("scanner", "structure")
workflow.add_edge("structure", "content")
workflow.add_edge("content", "reviewer")

# 조건부 연결 (Loop)
workflow.add_conditional_edges(
    "reviewer",
    route_after_review,
    {
        "content": "content",
        "renderer": "renderer"
    }
)

workflow.add_edge("renderer", END)

# 컴파일
app = workflow.compile()


from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt

# [Helper] 더러운 데이터를 숫자로 씻어주는 함수
def sanitize_number(value):
    if isinstance(value, (int, float)):
        return value
    try:
        # "1,000" -> 1000.0, "10%" -> 10.0 처리 등을 여기서 함
        clean_str = str(value).replace(",", "").replace("%", "").strip()
        return float(clean_str)
    except:
        return 0.0 # 정 안되면 0으로 처리

def draw_chart_safe(slide, x, y, w, h, data_dict):
    try:
        # 1. 데이터 꺼내기 (Pydantic 모델이 dict로 변환되어 들어옴)
        labels = data_dict.get("chart_labels", []) or []
        raw_values = data_dict.get("chart_values", []) or []
        title = data_dict.get("chart_title", "")

        # 2. 데이터 유효성 검사 (데이터 없으면 그리기 중단)
        if not labels or not raw_values:
            print(f"   ⚠️ 차트 데이터 누락 (Labels: {len(labels)}, Values: {len(raw_values)})")
            return

        # 3. 값(Values) 안전하게 숫자로 변환
        values = [sanitize_number(v) for v in raw_values]

        # 4. [중요] X축과 Y축 개수 맞추기 (짧은 쪽에 맞춤)
        min_len = min(len(labels), len(values))
        labels = labels[:min_len]
        values = values[:min_len]

        # 5. 차트 데이터 객체 생성
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series(title or "Series 1", values)

        # 6. PPT에 삽입
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
        ).chart

        # 7. 제목 설정
        if title:
            chart.chart_title.text_frame.text = title
            
        print("   ✅ 차트 생성 성공")

    except Exception as e:
        print(f"   ❌ 차트 렌더링 에러: {e}")
        # 실패 시 빈 자리에 에러 메시지라도 남겨둠 (디버깅용)
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.text_frame.text = f"[Chart Error]\n{str(e)}"

def draw_table_safe(slide, x, y, w, h, data_dict):
    try:
        rows = data_dict.get("table_rows", [])
        if not rows: return

        # 행/열 개수 계산
        r_cnt = len(rows)
        c_cnt = max(len(r) for r in rows) if r_cnt > 0 else 0
        
        if r_cnt == 0 or c_cnt == 0: return

        # 테이블 생성
        graphic_frame = slide.shapes.add_table(r_cnt, c_cnt, x, y, w, h)
        table = graphic_frame.table

        # 셀 채우기
        for i, row_data in enumerate(rows):
            for j, cell_val in enumerate(row_data):
                # 데이터가 짧아서 인덱스 에러나는 것 방지
                if j >= c_cnt: break 
                
                cell = table.cell(i, j)
                cell.text = str(cell_val)
                # (옵션) 폰트 사이즈 조정
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        print("   ✅ 테이블 생성 성공")
        
    except Exception as e:
        print(f"   ❌ 테이블 렌더링 에러: {e}")
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.text_frame.text = f"[Table Error]\n{str(e)}"


# renderer.py (또는 헬퍼 함수 정의 부분)

from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# [Config] 차트 타입 매핑 사전
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA
}

# [Config] 테이블 스타일 ID 매핑 (PPT 내부 GUID)
# 자주 쓰는 스타일 몇 개만 매핑해두면 편합니다.
TABLE_STYLE_MAP = {
    "light": "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",   # Medium Style 2 - Accent 1
    "medium": "{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}",  # (위와 동일, 취향껏 변경 가능)
    "dark": "{2D5ABB26-0587-4C30-8999-92F81FD0307C}",    # Themed Style 1 - Accent 1
    "accent": "{3C2FFA5D-87B4-456A-9821-1D502468CF0F}"   # Medium Style 4 - Accent 1
}

def draw_chart_safe(slide, x, y, w, h, data_dict):
    try:
        # 1. 데이터 파싱 (기존 로직 동일)
        labels = data_dict.get("chart_labels", []) or []
        raw_values = data_dict.get("chart_values", []) or []
        title = data_dict.get("chart_title", "")
        # [NEW] 차트 타입 가져오기
        c_type_str = data_dict.get("chart_type", "bar").lower()
        
        # ... (중간 데이터 정제 로직 sanitize_number 등은 기존 유지) ...
        values = [sanitize_number(v) for v in raw_values] # (예시)

        # 2. 차트 데이터 객체 생성
        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series(title or "Series 1", values)

        # 3. [NEW] 선택된 차트 타입으로 그리기
        ppt_chart_type = CHART_TYPE_MAP.get(c_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        chart = slide.shapes.add_chart(
            ppt_chart_type, x, y, w, h, chart_data
        ).chart

        # 4. 옵션: 차트 종류별 미세 조정 (예: 원형 차트는 범례가 중요)
        if c_type_str in ["pie", "doughnut"]:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # 5. 제목 설정
        if title:
            chart.chart_title.text_frame.text = title
            
        print(f"   ✅ 차트 생성 성공 ({c_type_str})")

    except Exception as e:
        print(f"   ❌ 차트 렌더링 에러: {e}")
        # ... (에러 처리 로직) ...

def draw_table_safe(slide, x, y, w, h, data_dict):
    try:
        rows = data_dict.get("table_rows", [])
        # [NEW] 스타일 가져오기
        style_key = data_dict.get("table_style", "medium") 
        
        if not rows: return
        r_cnt, c_cnt = len(rows), len(rows[0])

        graphic_frame = slide.shapes.add_table(r_cnt, c_cnt, x, y, w, h)
        table = graphic_frame.table

        # [NEW] 테이블 스타일 적용
        # python-pptx는 table_style_id에 GUID 문자열을 넣어야 합니다.
        target_style_id = TABLE_STYLE_MAP.get(style_key, TABLE_STYLE_MAP["medium"])
        table.table_style_id = target_style_id

        # ... (셀 채우기 로직 기존 유지) ...
        
        print(f"   ✅ 테이블 생성 성공 (Style: {style_key})")
        
    except Exception as e:
        print(f"   ❌ 테이블 렌더링 에러: {e}")


def replace_text_preserving_style(shape, new_text):
    """
    [핵심 기능]
    기존 텍스트 상자의 폰트, 색상, 크기, 볼드체 등을 그대로 유지하면서
    글자 내용만 'new_text'로 싹 바꿔치기합니다.
    """
    if not shape.has_text_frame: 
        return
    
    tf = shape.text_frame
    # 기존에 글자가 없으면 그냥 넣고 끝냄
    if not tf.paragraphs:
        tf.text = new_text
        return

    # 1. 첫 번째 문단의 첫 번째 스타일(Run)을 '샘플'로 복사
    p = tf.paragraphs[0]
    sample_run = p.runs[0] if p.runs else None
    
    # 스타일 백업
    font_name = sample_run.font.name if sample_run else None
    font_size = sample_run.font.size if sample_run else None
    font_color = sample_run.font.color.rgb if (sample_run and hasattr(sample_run.font.color, 'rgb')) else None
    is_bold = sample_run.font.bold if sample_run else None

    # 2. 내용 교체 (기존 것 다 지움)
    tf.clear() 
  

from pptx.enum.shapes import PP_PLACEHOLDER

def smart_fill_placeholders(slide, data_dict):
    """
    data_dict: {"Title": "...", "Subtitle": "...", "Content": "..."}
    """
    # 에이전트가 준 키들을 전부 소문자로 바꿔서 검색하기 쉽게 만듦
    # 예: {"title": "...", "main_title": "..."}
    normalized_data = {k.lower(): v for k, v in data_dict.items()}

    for shape in slide.placeholders:
        # 1. 실제 PPT의 Placeholder 정보 확인
        ph_type = shape.placeholder_format.type
        ph_name = shape.name.lower()
        
        target_text = None

        # --- [매칭 로직 1] 타입(Type)으로 찾기 (가장 정확함) ---
        # (1) 제목 칸 (CENTER_TITLE or TITLE)
        if ph_type == PP_PLACEHOLDER.CENTER_TITLE or ph_type == PP_PLACEHOLDER.TITLE:
            # 에이전트가 title, main_title, subject 중 하나라도 보냈으면 씀
            target_text = normalized_data.get("title") or normalized_data.get("main_title") or normalized_data.get("subject")

        # (2) 부제목 칸 (SUBTITLE)
        elif ph_type == PP_PLACEHOLDER.SUBTITLE:
            target_text = normalized_data.get("subtitle") or normalized_data.get("sub_title")

        # (3) 본문/바디 칸 (BODY or OBJECT)
        elif ph_type == PP_PLACEHOLDER.BODY or ph_type == PP_PLACEHOLDER.OBJECT:
            # content, body, description 중 하나라도 보냈으면 씀
            target_text = normalized_data.get("content") or normalized_data.get("body") or normalized_data.get("desc")

        # (4) 날짜 (DATE)
        elif ph_type == PP_PLACEHOLDER.DATE:
            target_text = normalized_data.get("date")

        # (5) 쪽번호 (SLIDE_NUMBER) -> 보통 자동이지만 강제 입력 원할 때
        elif ph_type == PP_PLACEHOLDER.SLIDE_NUMBER:
            target_text = normalized_data.get("page_no")

        # --- [매칭 로직 2] 이름(Name)으로 찾기 (타입 매칭 실패 시) ---
        if not target_text:
            # PPT 이름이 "Content Placeholder 2"라면 -> "content"라는 키가 있는지 확인
            for key, val in normalized_data.items():
                if key in ph_name: # 부분 일치 검색
                    target_text = val
                    break
        
        # 2. 찾았으면 갈아끼우기 (여기서 replace 함수 사용!)
        if target_text:
            replace_text_preserving_style(shape, target_text)
            print(f"   ✅ Placeholder 채움: {shape.name} <- '{target_text[:10]}...'")
        else:
            # 디버깅용 로그: 왜 안 들어갔는지 확인 가능
            print(f"   ⚠️ 매칭 실패: PPT칸({shape.name}/{ph_type}) vs 데이터키({list(normalized_data.keys())})")


def renderer_node(state: AgentState):
    # ... (상단 생략) ...
    
    for plan in state["slide_data"]:
        # ... (슬라이드 생성) ...
        
        # [기존 코드 삭제]
        # common = plan.get("common_fields", {})
        # for shape in slide.placeholders:
        #     ... (복잡했던 if/else 로직) ...

        # [NEW: 한 줄로 끝내기]
        # 1. 공통 필드(제목, 본문 등) 채우기
        common_data = plan.get("common_fields", {})
        smart_fill_placeholders(slide, common_data)
        
        # 2. Dynamic Components 그리기 (차트 등)
        # ... (이건 기존 유지) ...
  
    # 3. 새 내용 넣고 스타일 복원 (수술 완료)
    new_p = tf.paragraphs[0]
    new_run = new_p.add_run()
    new_run.text = str(new_text) # 안전하게 문자열 변환

    if sample_run:
        if font_name: new_run.font.name = font_name
        if font_size: new_run.font.size = font_size
        if font_color: new_run.font.color.rgb = font_color
        if is_bold is not None: new_run.font.bold = is_bold


# renderer.py 내부 혹은 utils.py

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_str):
    """ 'FFFFFF' -> RGBColor(255, 255, 255) 변환 """
    if not hex_str or len(hex_str) != 6: return None
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:], 16))

def draw_table_advanced(slide, x, y, w, h, data_dict):
    try:
        rows = data_dict.get("table_rows", [])
        if not rows: return

        r_cnt = len(rows)
        c_cnt = max(len(r) for r in rows) if r_cnt > 0 else 0
        if r_cnt == 0 or c_cnt == 0: return

        # 1. 테이블 생성
        graphic_frame = slide.shapes.add_table(r_cnt, c_cnt, x, y, w, h)
        table = graphic_frame.table

        # [핵심 1] 행 높이 강제 배분 (Bounding Box 준수)
        # 전체 높이 h를 행 개수로 나누어 할당
        row_height = h / r_cnt
        for row in table.rows:
            row.height = int(row_height)

        # [핵심 2] 폰트 크기 자동 계산 (Auto-Sizing Logic)
        # 사용자가 지정했으면 그거 쓰고, 아니면 행 개수에 따라 작게 조절
        user_font_size = data_dict.get("table_font_size")
        
        if user_font_size:
            final_font_size = Pt(user_font_size)
        else:
            # 휴리스틱: 행이 10개 넘으면 10pt, 5개 넘으면 12pt, 아니면 14pt
            if r_cnt > 15: final_font_size = Pt(9)
            elif r_cnt > 10: final_font_size = Pt(10)
            elif r_cnt > 5: final_font_size = Pt(12)
            else: final_font_size = Pt(14)

        # 3. 셀 데이터 채우기 및 스타일 적용
        for i, row_data in enumerate(rows):
            for j, cell_val in enumerate(row_data):
                if j >= len(table.columns): break
                
                cell = table.cell(i, j)
                
                # 텍스트가 셀을 넘치지 않게 여백 줄이기 (선택사항)
                cell.margin_left = Pt(2)
                cell.margin_right = Pt(2)
                cell.margin_top = Pt(1)
                cell.margin_bottom = Pt(1)
                
                # 수직 정렬 (중앙)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # 텍스트 입력
                # (기존 텍스트 프레임 초기화 후 입력해야 서식 적용이 깔끔함)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(cell_val)
                
                # [폰트 적용]
                run.font.size = final_font_size
                run.font.name = "맑은 고딕" # 또는 "Arial" 등 원하는 폰트
                
                # (옵션) 헤더(첫 줄) 스타일링
                if i == 0:
                    run.font.bold = True
                    # 사용자 지정 헤더 색상
                    header_color = data_dict.get("table_header_color")
                    if header_color:
                        try:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = hex_to_rgb(header_color)
                        except: pass # 색상 코드 오류 시 무시

        # [NEW] 스타일 ID 적용 (배경색 등 기본 테마)
        style_key = data_dict.get("table_style", "medium")
        # table.table_style_id = TABLE_STYLE_MAP.get(...) # (이전 단계 코드 활용)
        
        print(f"   ✅ 테이블 생성 (Rows: {r_cnt}, FontSize: {final_font_size.pt}pt)")

    except Exception as e:
        print(f"   ❌ 테이블 렌더링 에러: {e}")


# schema.py

from typing import List, Optional, Literal
from pydantic import BaseModel, Field

# [NEW] 1. 하위 모델 정의: 하나의 데이터 시리즈 (예: '2023년 매출'과 그 값들)
class ChartSeries(BaseModel):
    name: str = Field(..., description="범례(Legend)에 표시될 시리즈 이름 (예: '영업이익', '순이익')")
    values: List[float] = Field(..., description="해당 시리즈의 데이터 값 리스트 (숫자만)")

# 2. 메인 데이터 모델 수정
class ComponentData(BaseModel):
    # ... (text_content, table 관련 필드 유지) ...

    chart_title: Optional[str] = Field(None, description="차트 제목")
    chart_labels: Optional[List[str]] = Field(None, description="X축 라벨 리스트 (모든 시리즈 공통)")
    
    # 🚨 [핵심 수정] 기존 chart_values 필드를 삭제하고 아래로 대체
    chart_series: Optional[List[ChartSeries]] = Field(
        None, 
        description="다중 시리즈 데이터. 꺾은선 2개 이상, 묶은 세로 막대형 등 복합 차트 구현 시 사용."
    )
    
    chart_type: Literal["bar", "line", "pie", "doughnut", "area"] = Field(
        "bar", 
        description="차트 종류 (bar: 묶은 세로 막대, line: 꺾은선)"
    )


# renderer.py

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# (CHART_TYPE_MAP, sanitize_number 함수는 기존과 동일하다고 가정)

def draw_chart_advanced(slide, x, y, w, h, data_dict):
    try:
        # 1. 공통 데이터 추출
        title = data_dict.get("chart_title", "")
        labels = data_dict.get("chart_labels", []) or []
        c_type_str = data_dict.get("chart_type", "bar").lower()
        
        # 🚨 [핵심 1] 다중 시리즈 데이터 추출 (Pydantic 모델 -> dict 리스트로 변환됨)
        raw_series_list = data_dict.get("chart_series", []) or []

        # 데이터 유효성 검사
        if not labels or not raw_series_list:
            print("   ⚠️ 차트 데이터 누락 (Labels 또는 Series 없음)")
            return

        # 2. 차트 데이터 객체 생성 및 라벨 설정
        chart_data = CategoryChartData()
        chart_data.categories = labels
        
        # 🚨 [핵심 2] 반복문을 돌며 시리즈 추가 (Multi-Series Logic)
        label_len = len(labels)
        for series in raw_series_list:
            s_name = series.get("name", "Series")
            s_values_raw = series.get("values", [])
            
            # 값 정제 (숫자 변환)
            s_values_clean = [sanitize_number(v) for v in s_values_raw]
            
            # 길이 맞춤 (라벨 개수만큼 잘라내기)
            s_values_final = s_values_clean[:label_len]
            
            # 데이터 추가
            chart_data.add_series(s_name, s_values_final)
            print(f"      + 시리즈 추가: {s_name} (데이터 {len(s_values_final)}개)")

        # 3. 차트 생성
        ppt_chart_type = CHART_TYPE_MAP.get(c_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart = slide.shapes.add_chart(
            ppt_chart_type, x, y, w, h, chart_data
        ).chart

        # 4. 제목 및 범례(Legend) 설정
        if title:
            chart.chart_title.text_frame.text = title
        
        # 시리즈가 2개 이상이거나 파이 차트면 범례 표시
        if len(raw_series_list) > 1 or c_type_str in ["pie", "doughnut"]:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        print(f"   ✅ 고도화된 차트 생성 성공 ({c_type_str}, 시리즈 {len(raw_series_list)}개)")

    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"   ❌ 차트 렌더링 에러: {e}")
        # (에러 시 텍스트 박스 대체 로직 유지)






